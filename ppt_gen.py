import os
import tempfile
import streamlit as st
from PyPDF2 import PdfReader
import pdfplumber
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from typing import Optional, Tuple
from dotenv import load_dotenv
import re

# ========== CONFIGURATION ========== #
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

MODEL_NAME = "gemini-2.0-flash"
MAX_PDF_SIZE_MB = 50
MAX_SLIDES = 10
PROCESSING_CHUNK_SIZE = 15000
MIN_CONTENT_LENGTH = 100
PPT_FONT = "Calibri"

def configure_gemini() -> genai.GenerativeModel:
    try:
        return genai.GenerativeModel(
            MODEL_NAME,
            generation_config={
                "temperature": 0.3,
                "top_p": 0.95,
                "top_k": 40,
                "max_output_tokens": 2048,
            }
        )
    except Exception as e:
        st.error(f"Failed to configure Gemini: {str(e)}")
        st.stop()

def extract_text_from_pdf(pdf_file) -> Tuple[Optional[str], Optional[str]]:
    pdf_file.seek(0)
    text = ""
    try:
        with pdfplumber.open(pdf_file) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n\n[Page {i+1}]\n{page_text}"
                if len(text) > PROCESSING_CHUNK_SIZE:
                    break
    except Exception as e:
        st.warning(f"pdfplumber failed: {str(e)}")

    if len(text) < MIN_CONTENT_LENGTH:
        pdf_file.seek(0)
        try:
            reader = PdfReader(pdf_file)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n\n[Page {i+1}]\n{page_text}"
                if len(text) > PROCESSING_CHUNK_SIZE:
                    break
        except Exception as e:
            st.warning(f"PyPDF2 failed: {str(e)}")

    if len(text) > MIN_CONTENT_LENGTH:
        return text, None
    return None, "Failed to extract sufficient text (document may be scanned)"

def generate_slide_structure(model, pdf_text: str, ppt_title: str) -> Tuple[Optional[str], Optional[str]]:
    try:
        prompt = f"""
You are an expert presentation designer. Based on the content below, create a PowerPoint structure for the title: "{ppt_title}".

PDF CONTENT:
{pdf_text[:PROCESSING_CHUNK_SIZE]}

Return exactly 5 slides in this format:
**Slide 1: [Title Slide]**
* **Title:** "{ppt_title}"
* **Subtitle:** "[1-line summary]"

**Slide 2: [Introduction]**
* **Title:** "Intro title"
* **Bullet Points:**
    * Bullet 1
    * Bullet 2

...repeat for 5 slides total.
"""
        response = model.generate_content(prompt)
        return response.text.strip(), None
    except Exception as e:
        return None, f"Structure generation failed: {str(e)}"

def parse_gemini_structure(slide_structure: str):
    slides = re.split(r'\*\*Slide \d+:.*?\*\*', slide_structure)
    headers = re.findall(r'\*\*Slide \d+:.*?\*\*', slide_structure)

    parsed = []
    for i, header in enumerate(headers):
        content = slides[i+1] if i+1 < len(slides) else ""
        title_match = re.search(r'\*\*Title:\*\*\s*(.*)', content)
        title = title_match.group(1).strip() if title_match else f"Slide {i+2}"
        bullets = re.findall(r'\*\s+(.*)', content)
        parsed.append((title, bullets))
    return parsed

def create_presentation(ppt_title: str, slide_structure: str, model, pdf_text: str) -> Optional[str]:
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = ppt_title
        slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%d %B %Y')}"

        parsed_slides = parse_gemini_structure(slide_structure)

        if not parsed_slides:
            st.warning("⚠️ Could not parse any slides from Gemini output. Using fallback.")
            parsed_slides = [("Introduction", ["Key points missing from AI output"])]

        for title, bullet_points in parsed_slides[:MAX_SLIDES]:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]

            title_shape.text = title
            text_frame = content_shape.text_frame
            text_frame.clear()

            for i, bullet in enumerate(bullet_points):
                para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                para.text = bullet.strip("-•* ").strip()
                para.font.size = Pt(18)
                para.font.name = PPT_FONT
                para.font.color.rgb = RGBColor(0, 0, 0)
                para.alignment = PP_ALIGN.LEFT

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs.save(tmp.name)
            return tmp.name
    except Exception as e:
        st.error(f"PPT creation failed: {str(e)}")
        return None

def main():
    st.set_page_config(page_title="PDF to PowerPoint Pro", page_icon="📊", layout="centered")
    st.title("📊 PDF to PowerPoint Converter")
    st.write("Upload a PDF to generate a professional PowerPoint presentation")

    model = configure_gemini()
    pdf_file = st.file_uploader("Choose PDF file", type=["pdf"])

    if pdf_file:
        if pdf_file.size > MAX_PDF_SIZE_MB * 1024 * 1024:
            st.error(f"File too large (max {MAX_PDF_SIZE_MB}MB)")
            return

        with st.spinner("Extracting text from PDF..."):
            pdf_text, error = extract_text_from_pdf(pdf_file)
            if error:
                st.error(error)
                return
            st.session_state.pdf_text = pdf_text

        ppt_title = st.text_input("Presentation Title", "Business Report")

        if st.button("Analyze Document"):
            with st.spinner("Creating slide structure..."):
                slide_structure, error = generate_slide_structure(model, pdf_text, ppt_title)
                if error:
                    st.error(error)
                else:
                    st.session_state.slide_structure = slide_structure
                    st.subheader("📋 Slide Structure (Preview)")
                    st.code(slide_structure)

        if 'slide_structure' in st.session_state:
            edited_structure = st.text_area("✏️ Edit structure if needed:", value=st.session_state.slide_structure, height=300)

            if st.button("Generate PowerPoint", type="primary"):
                with st.spinner("Creating presentation..."):
                    ppt_path = create_presentation(ppt_title, edited_structure, model, st.session_state.pdf_text)
                    if ppt_path:
                        with open(ppt_path, "rb") as f:
                            ppt_bytes = f.read()
                        st.success("✅ Presentation generated successfully!")
                        st.download_button(
                            label="📥 Download PowerPoint",
                            data=ppt_bytes,
                            file_name=f"{ppt_title.replace(' ', '_')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        try:
                            os.unlink(ppt_path)
                        except:
                            pass

if __name__ == "__main__":
    main()
